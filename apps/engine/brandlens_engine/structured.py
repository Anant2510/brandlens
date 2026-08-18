"""Structured-source parsing: PDF / PPTX / HTML / Figma JSON -> one element tree.

Rule 1 of the engine: **prefer structured sources over pixels.** If the asset
carries its own description of itself, an inferred measurement is strictly worse
than a read one. OCR guesses "Helvetica, about 11pt"; a PDF font descriptor
*states* `ABCDEF+HelveticaNeue-Bold, 10.5pt, fill #1A1A1A` with an exact bbox.
Every downstream analyzer therefore consumes `StructuredDocument` when one is
available and falls back to pixels only when it is not — but both paths converge
on the same `CriterionResult` schema so a finding reads identically either way.
"""

from __future__ import annotations

import io
import re
from dataclasses import dataclass, field
from html.parser import HTMLParser
from typing import Any, Literal

from .logging import get_logger

log = get_logger(__name__)

SourceKind = Literal["pdf", "pptx", "html", "figma", "none"]

# Weight tokens that indicate a *real* bold face exists in the font name. When
# a run declares bold and none of these appear, the renderer is smearing the
# regular face — faux bold, which brand teams treat as a typography violation.
_BOLD_TOKENS = (
    "bold", "black", "heavy", "semibold", "demibold", "extrabold", "ultrabold",
    "medium", "600", "700", "800", "900",
)
_ITALIC_TOKENS = ("italic", "oblique", "it", "slanted")
_SUBSET_PREFIX = re.compile(r"^[A-Z]{6}\+")


@dataclass(slots=True)
class TextElement:
    text: str
    #: Normalized to the page, origin top-left.
    bbox: tuple[float, float, float, float]
    font_family: str = ""
    font_size_pt: float = 0.0
    font_weight: int = 400
    is_bold: bool = False
    is_italic: bool = False
    #: Declared bold/italic with no matching face in the font name.
    is_faux_bold: bool = False
    is_faux_italic: bool = False
    color_hex: str | None = None
    page: int = 0
    ordinal: int = 0
    line_height_pt: float | None = None
    embedded: bool = True
    role_hint: str | None = None

    @property
    def height_norm(self) -> float:
        return max(0.0, self.bbox[3] - self.bbox[1])


@dataclass(slots=True)
class ShapeElement:
    bbox: tuple[float, float, float, float]
    fill_hex: str | None = None
    stroke_hex: str | None = None
    kind: str = "rect"
    page: int = 0
    area_norm: float = 0.0


@dataclass(slots=True)
class ImageElement:
    bbox: tuple[float, float, float, float]
    page: int = 0
    name: str | None = None
    xref: int | None = None


@dataclass(slots=True)
class StructuredPage:
    index: int
    width_pt: float
    height_pt: float
    text: list[TextElement] = field(default_factory=list)
    shapes: list[ShapeElement] = field(default_factory=list)
    images: list[ImageElement] = field(default_factory=list)

    @property
    def plain_text(self) -> str:
        return "\n".join(t.text for t in self.text if t.text.strip())


@dataclass(slots=True)
class StructuredDocument:
    kind: SourceKind
    pages: list[StructuredPage] = field(default_factory=list)
    #: family -> embedded?  A non-embedded family is a live substitution risk.
    fonts: dict[str, bool] = field(default_factory=dict)
    warnings: list[str] = field(default_factory=list)
    meta: dict[str, Any] = field(default_factory=dict)

    @property
    def available(self) -> bool:
        return self.kind != "none" and bool(self.pages)

    @property
    def all_text(self) -> list[TextElement]:
        return [t for p in self.pages for t in p.text]

    @property
    def plain_text(self) -> str:
        return "\n".join(p.plain_text for p in self.pages)

    def page(self, index: int = 0) -> StructuredPage | None:
        return self.pages[index] if 0 <= index < len(self.pages) else None


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------
def normalize_font_name(name: str) -> str:
    """Strip PDF subset prefixes and style suffixes down to a family name."""
    n = _SUBSET_PREFIX.sub("", (name or "").strip())
    n = n.replace("_", "-")
    parts = re.split(r"[-,]", n)
    family = parts[0] if parts else n
    family = re.sub(r"(?i)(MT|PS|Std|Pro|LT)$", "", family).strip()
    return re.sub(r"(?<=[a-z])(?=[A-Z])", " ", family).strip()


def _name_declares(name: str, tokens: tuple[str, ...]) -> bool:
    low = _SUBSET_PREFIX.sub("", (name or "")).lower()
    return any(t in low for t in tokens)


def _detect_faux(raw_font: str, declared_bold: bool, declared_italic: bool) -> tuple[bool, bool]:
    faux_bold = declared_bold and not _name_declares(raw_font, _BOLD_TOKENS)
    faux_italic = declared_italic and not _name_declares(raw_font, _ITALIC_TOKENS)
    return faux_bold, faux_italic


def _weight_from_name(name: str, declared_bold: bool) -> int:
    low = (name or "").lower()
    for token, weight in (
        ("thin", 100), ("extralight", 200), ("ultralight", 200), ("light", 300),
        ("regular", 400), ("book", 400), ("medium", 500), ("semibold", 600),
        ("demibold", 600), ("extrabold", 800), ("ultrabold", 800), ("black", 900),
        ("heavy", 900), ("bold", 700),
    ):
        if token in low:
            return weight
    return 700 if declared_bold else 400


def _int_to_hex(value: int) -> str:
    v = int(value) & 0xFFFFFF
    return f"#{(v >> 16) & 255:02X}{(v >> 8) & 255:02X}{v & 255:02X}"


def _floats_to_hex(rgb: object) -> str | None:
    if not isinstance(rgb, (list, tuple)) or len(rgb) < 3:
        return None
    try:
        r, g, b = (max(0.0, min(1.0, float(c))) for c in tuple(rgb)[:3])
    except (TypeError, ValueError):
        return None
    return f"#{int(round(r * 255)):02X}{int(round(g * 255)):02X}{int(round(b * 255)):02X}"


def _norm(box: tuple[float, float, float, float], w: float, h: float) -> tuple[float, float, float, float]:
    w = max(w, 1e-6)
    h = max(h, 1e-6)
    clamp = lambda v: float(min(1.0, max(0.0, v)))  # noqa: E731
    return (clamp(box[0] / w), clamp(box[1] / h), clamp(box[2] / w), clamp(box[3] / h))


# ---------------------------------------------------------------------------
# PDF (PyMuPDF)
# ---------------------------------------------------------------------------
def parse_pdf(data: bytes, max_pages: int = 40) -> StructuredDocument:
    """Exact fonts/sizes/colours/bboxes from the PDF itself.

    `get_text("dict")` gives per-span font descriptor data; `get_drawings()`
    gives vector fills, which is how we find the flat brand-colour panels that
    a pixel palette would blur into their neighbours.
    """
    doc = StructuredDocument(kind="pdf")
    try:
        import pymupdf as fitz  # `fitz` is the legacy alias; import the new name
    except ImportError:
        doc.kind = "none"
        doc.warnings.append("PyMuPDF unavailable; falling back to pixel analysis")
        return doc

    try:
        pdf = fitz.open(stream=data, filetype="pdf")
    except Exception as exc:  # noqa: BLE001 - a corrupt file degrades, never 500s
        doc.kind = "none"
        doc.warnings.append(f"unreadable PDF: {exc}")
        return doc

    try:
        if pdf.is_encrypted and not pdf.authenticate(""):
            doc.kind = "none"
            doc.warnings.append("PDF is encrypted; structured extraction unavailable")
            return doc

        doc.meta["page_count"] = pdf.page_count
        for pno in range(min(pdf.page_count, max_pages)):
            try:
                page = pdf[pno]
            except Exception as exc:  # noqa: BLE001
                doc.warnings.append(f"page {pno} unreadable: {exc}")
                continue

            rect = page.rect
            pw, ph = float(rect.width), float(rect.height)
            spage = StructuredPage(index=pno, width_pt=pw, height_pt=ph)

            for xref, ext, _ftype, basefont, *_rest in page.get_fonts(full=True):
                fam = normalize_font_name(basefont)
                embedded = str(ext).lower() not in ("", "n/a", "none")
                if fam:
                    doc.fonts[fam] = doc.fonts.get(fam, False) or embedded
                del xref

            ordinal = 0
            try:
                td = page.get_text("dict")
            except Exception as exc:  # noqa: BLE001
                doc.warnings.append(f"page {pno} text extraction failed: {exc}")
                td = {"blocks": []}

            for block in td.get("blocks", []):
                if block.get("type") == 1:  # image block
                    bb = block.get("bbox")
                    if bb:
                        spage.images.append(
                            ImageElement(bbox=_norm(tuple(bb), pw, ph), page=pno, xref=block.get("number"))
                        )
                    continue
                for line in block.get("lines", []):
                    line_bbox = line.get("bbox")
                    for span in line.get("spans", []):
                        text = span.get("text", "")
                        if not text.strip():
                            continue
                        raw_font = str(span.get("font", ""))
                        flags = int(span.get("flags", 0))
                        declared_bold = bool(flags & 16)
                        declared_italic = bool(flags & 2)
                        faux_b, faux_i = _detect_faux(raw_font, declared_bold, declared_italic)
                        family = normalize_font_name(raw_font)
                        bbox = tuple(span.get("bbox", line_bbox or (0, 0, 0, 0)))
                        lh = None
                        if line_bbox:
                            lh = float(line_bbox[3]) - float(line_bbox[1])
                        spage.text.append(
                            TextElement(
                                text=text,
                                bbox=_norm(bbox, pw, ph),  # type: ignore[arg-type]
                                font_family=family,
                                font_size_pt=float(span.get("size", 0.0)),
                                font_weight=_weight_from_name(raw_font, declared_bold),
                                is_bold=declared_bold,
                                is_italic=declared_italic,
                                is_faux_bold=faux_b,
                                is_faux_italic=faux_i,
                                color_hex=_int_to_hex(int(span.get("color", 0))),
                                page=pno,
                                ordinal=ordinal,
                                line_height_pt=lh,
                                embedded=doc.fonts.get(family, True),
                            )
                        )
                        ordinal += 1

            try:
                for d in page.get_drawings():
                    r = d.get("rect")
                    if r is None:
                        continue
                    box = (float(r.x0), float(r.y0), float(r.x1), float(r.y1))
                    area = ((box[2] - box[0]) * (box[3] - box[1])) / max(pw * ph, 1e-6)
                    spage.shapes.append(
                        ShapeElement(
                            bbox=_norm(box, pw, ph),
                            fill_hex=_floats_to_hex(d.get("fill")),
                            stroke_hex=_floats_to_hex(d.get("color")),
                            kind=str(d.get("type", "rect")),
                            page=pno,
                            area_norm=float(area),
                        )
                    )
            except Exception as exc:  # noqa: BLE001
                doc.warnings.append(f"page {pno} vector extraction failed: {exc}")

            doc.pages.append(spage)
    finally:
        try:
            pdf.close()
        except Exception:  # noqa: BLE001,S110  # pragma: no cover
            pass

    if not doc.pages:
        doc.kind = "none"
        doc.warnings.append("PDF produced no readable pages")
    return doc


# ---------------------------------------------------------------------------
# PPTX (python-pptx)
# ---------------------------------------------------------------------------
_EMU_PER_PT = 12700.0


def parse_pptx(data: bytes, max_pages: int = 60) -> StructuredDocument:
    doc = StructuredDocument(kind="pptx")
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        doc.kind = "none"
        doc.warnings.append("python-pptx unavailable")
        return doc

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        doc.kind = "none"
        doc.warnings.append(f"unreadable PPTX: {exc}")
        return doc

    sw = float(Emu(prs.slide_width).pt) if prs.slide_width else 720.0
    sh = float(Emu(prs.slide_height).pt) if prs.slide_height else 540.0

    for idx, slide in enumerate(prs.slides):
        if idx >= max_pages:
            break
        spage = StructuredPage(index=idx, width_pt=sw, height_pt=sh)
        ordinal = 0
        for shape in slide.shapes:
            try:
                left = float(Emu(shape.left or 0).pt)
                top = float(Emu(shape.top or 0).pt)
                width = float(Emu(shape.width or 0).pt)
                height = float(Emu(shape.height or 0).pt)
            except (TypeError, ValueError):
                continue
            box = (left, top, left + width, top + height)
            nbox = _norm(box, sw, sh)

            if shape.shape_type is not None and "PICTURE" in str(shape.shape_type):
                spage.images.append(ImageElement(bbox=nbox, page=idx, name=getattr(shape, "name", None)))

            fill_hex: str | None = None
            try:
                fill = shape.fill
                if fill.type is not None and str(fill.type).startswith("SOLID"):
                    fill_hex = f"#{fill.fore_color.rgb!s}"
            except Exception:  # noqa: BLE001 - theme/gradient fills raise; not an error
                fill_hex = None
            if fill_hex or not shape.has_text_frame:
                spage.shapes.append(
                    ShapeElement(
                        bbox=nbox,
                        fill_hex=fill_hex,
                        kind=str(shape.shape_type),
                        page=idx,
                        area_norm=(width * height) / max(sw * sh, 1e-6),
                    )
                )

            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    f = run.font
                    raw_font = f.name or ""
                    declared_bold = bool(f.bold)
                    declared_italic = bool(f.italic)
                    size_pt = float(f.size.pt) if f.size is not None else 18.0
                    color_hex = None
                    try:
                        if f.color is not None and f.color.type is not None:
                            color_hex = f"#{f.color.rgb!s}"
                    except Exception:  # noqa: BLE001 - theme colours have no .rgb
                        color_hex = None
                    faux_b, faux_i = _detect_faux(raw_font, declared_bold, declared_italic)
                    family = normalize_font_name(raw_font)
                    if family:
                        doc.fonts.setdefault(family, True)
                    spage.text.append(
                        TextElement(
                            text=run.text,
                            bbox=nbox,
                            font_family=family,
                            font_size_pt=size_pt,
                            font_weight=_weight_from_name(raw_font, declared_bold),
                            is_bold=declared_bold,
                            is_italic=declared_italic,
                            # PPTX bold is a *request* to the renderer; PowerPoint
                            # synthesises the face when the family has no bold cut.
                            is_faux_bold=faux_b,
                            is_faux_italic=faux_i,
                            color_hex=color_hex,
                            page=idx,
                            ordinal=ordinal,
                            role_hint=str(para.level) if para.level else None,
                        )
                    )
                    ordinal += 1
        doc.pages.append(spage)

    if not doc.pages:
        doc.kind = "none"
        doc.warnings.append("PPTX produced no slides")
    return doc


# ---------------------------------------------------------------------------
# HTML (stdlib parser — no extra dependency)
# ---------------------------------------------------------------------------
_CSS_SIZE = re.compile(r"font-size\s*:\s*([\d.]+)\s*(px|pt|em|rem|%)?", re.I)
_CSS_FAMILY = re.compile(r"font-family\s*:\s*([^;]+)", re.I)
_CSS_WEIGHT = re.compile(r"font-weight\s*:\s*([a-z0-9]+)", re.I)
_CSS_COLOR = re.compile(r"(?<!background-)color\s*:\s*(#[0-9a-f]{3,8}|rgba?\([^)]*\))", re.I)
_HEADING_PT = {"h1": 32.0, "h2": 24.0, "h3": 18.7, "h4": 16.0, "h5": 13.3, "h6": 10.7}


def _css_color_to_hex(value: str) -> str | None:
    v = value.strip()
    if v.startswith("#"):
        h = v.lstrip("#")
        if len(h) == 3:
            h = "".join(c * 2 for c in h)
        return f"#{h[:6].upper()}" if len(h) >= 6 else None
    m = re.findall(r"[\d.]+", v)
    if len(m) >= 3:
        try:
            r, g, b = (int(float(x)) for x in m[:3])
            return f"#{r:02X}{g:02X}{b:02X}"
        except ValueError:
            return None
    return None


class _HtmlCollector(HTMLParser):
    """Collects text runs with the inline style in force at that point.

    We do not lay out the DOM — there is no browser here — so bboxes are
    document-order estimates, and anything needing true geometry must come from
    a rendered screenshot instead. Fonts, sizes and colours, which is what the
    typography checks actually need, are exact when declared inline.
    """

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.runs: list[TextElement] = []
        self._stack: list[dict[str, Any]] = []
        self._skip = 0
        self._ordinal = 0

    def _current(self) -> dict[str, Any]:
        merged: dict[str, Any] = {"family": "", "size_pt": 12.0, "weight": 400, "italic": False, "color": None}
        for frame in self._stack:
            merged.update({k: v for k, v in frame.items() if v is not None})
        return merged

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in ("script", "style", "head"):
            self._skip += 1
            return
        style = ""
        for k, v in attrs:
            if k.lower() == "style" and v:
                style = v
        frame: dict[str, Any] = {"tag": tag}
        if tag in _HEADING_PT:
            frame["size_pt"] = _HEADING_PT[tag]
            frame["weight"] = 700
        if tag in ("b", "strong"):
            frame["weight"] = 700
        if tag in ("i", "em"):
            frame["italic"] = True
        if style:
            if (m := _CSS_FAMILY.search(style)) is not None:
                frame["family"] = m.group(1).split(",")[0].strip().strip("'\"")
            if (m := _CSS_SIZE.search(style)) is not None:
                val, unit = float(m.group(1)), (m.group(2) or "px").lower()
                frame["size_pt"] = {
                    "px": val * 0.75, "pt": val, "em": val * 12.0, "rem": val * 12.0, "%": val * 0.12
                }.get(unit, val * 0.75)
            if (m := _CSS_WEIGHT.search(style)) is not None:
                w = m.group(1).lower()
                frame["weight"] = 700 if w in ("bold", "bolder") else (int(w) if w.isdigit() else 400)
            if (m := _CSS_COLOR.search(style)) is not None:
                frame["color"] = _css_color_to_hex(m.group(1))
        self._stack.append(frame)

    def handle_endtag(self, tag: str) -> None:
        if tag in ("script", "style", "head"):
            self._skip = max(0, self._skip - 1)
            return
        for i in range(len(self._stack) - 1, -1, -1):
            if self._stack[i].get("tag") == tag:
                del self._stack[i:]
                break

    def handle_data(self, data: str) -> None:
        if self._skip or not data.strip():
            return
        st = self._current()
        weight = int(st["weight"])
        family = str(st["family"])
        declared_bold = weight >= 600
        faux_b, faux_i = _detect_faux(family, declared_bold, bool(st["italic"]))
        self.runs.append(
            TextElement(
                text=data.strip(),
                bbox=(0.0, 0.0, 1.0, 0.0),  # filled in by the caller in document order
                font_family=normalize_font_name(family),
                font_size_pt=float(st["size_pt"]),
                font_weight=weight,
                is_bold=declared_bold,
                is_italic=bool(st["italic"]),
                is_faux_bold=faux_b and bool(family),
                is_faux_italic=faux_i and bool(family),
                color_hex=st["color"],
                ordinal=self._ordinal,
                role_hint=str(self._stack[-1].get("tag")) if self._stack else None,
            )
        )
        self._ordinal += 1


def parse_html(text: str) -> StructuredDocument:
    doc = StructuredDocument(kind="html")
    collector = _HtmlCollector()
    try:
        collector.feed(text)
        collector.close()
    except Exception as exc:  # noqa: BLE001
        doc.kind = "none"
        doc.warnings.append(f"unreadable HTML: {exc}")
        return doc

    runs = collector.runs
    if not runs:
        doc.kind = "none"
        doc.warnings.append("HTML contained no text")
        return doc

    # Document-order vertical stacking: honest about being an estimate, and
    # good enough for reading order, which is all the copy checks need.
    n = len(runs)
    for i, r in enumerate(runs):
        r.bbox = (0.05, i / n, 0.95, (i + 1) / n)
        fam = r.font_family
        if fam:
            doc.fonts.setdefault(fam, False)  # web fonts are never "embedded"
    doc.warnings.append("HTML geometry is document-order estimated, not laid out")
    doc.pages.append(StructuredPage(index=0, width_pt=612.0, height_pt=792.0, text=runs))
    return doc


# ---------------------------------------------------------------------------
# Figma JSON
# ---------------------------------------------------------------------------
def parse_figma(payload: dict[str, Any]) -> StructuredDocument:
    """Walk a Figma REST `document` tree (or a single node payload)."""
    doc = StructuredDocument(kind="figma")
    root = payload.get("document") or payload.get("nodes") or payload
    if isinstance(root, dict) and "document" in root:
        root = root["document"]

    frames: list[dict[str, Any]] = []

    def _find_frames(node: object) -> None:
        if isinstance(node, dict):
            if node.get("type") in ("FRAME", "CANVAS", "COMPONENT", "PAGE") and node.get("absoluteBoundingBox"):
                frames.append(node)
            for child in node.get("children", []) or []:
                _find_frames(child)
            for v in node.values():
                if isinstance(v, dict) and "document" in v:
                    _find_frames(v["document"])
        elif isinstance(node, list):
            for item in node:
                _find_frames(item)

    _find_frames(root)
    if not frames and isinstance(root, dict) and root.get("absoluteBoundingBox"):
        frames = [root]
    if not frames:
        doc.kind = "none"
        doc.warnings.append("no Figma frames found in structuredSource")
        return doc

    for idx, frame in enumerate(frames[:40]):
        abb = frame.get("absoluteBoundingBox") or {}
        ox, oy = float(abb.get("x", 0.0)), float(abb.get("y", 0.0))
        fw, fh = float(abb.get("width", 1.0)) or 1.0, float(abb.get("height", 1.0)) or 1.0
        page = StructuredPage(index=idx, width_pt=fw, height_pt=fh)
        counter = [0]

        # Frame geometry is bound as defaults rather than closed over, so a
        # recursive walk cannot pick up the next frame's origin mid-traversal.
        def _walk(
            node: dict[str, Any],
            page: StructuredPage = page,
            counter: list[int] = counter,
            ox: float = ox,
            oy: float = oy,
            fw: float = fw,
            fh: float = fh,
        ) -> None:
            bb = node.get("absoluteBoundingBox")
            box = None
            if isinstance(bb, dict):
                x, y = float(bb.get("x", 0.0)) - ox, float(bb.get("y", 0.0)) - oy
                box = _norm((x, y, x + float(bb.get("width", 0.0)), y + float(bb.get("height", 0.0))), fw, fh)

            ntype = node.get("type")
            fills = node.get("fills") or []
            fill_hex = None
            for f in fills:
                if isinstance(f, dict) and f.get("type") == "SOLID" and f.get("visible", True):
                    c = f.get("color") or {}
                    fill_hex = _floats_to_hex((c.get("r"), c.get("g"), c.get("b")))
                    break

            if ntype == "TEXT" and box:
                style = node.get("style") or {}
                ps_name = str(style.get("fontPostScriptName") or "")
                raw_font = ps_name or str(style.get("fontFamily") or "")
                weight = int(style.get("fontWeight", 400) or 400)
                italic = bool(style.get("italic", False))
                # Only the PostScript name tells us which *face* Figma resolved.
                # `fontWeight` alone is the designer's request, so without the
                # PS name we cannot distinguish a real bold cut from a synthesised
                # one — and guessing here would fire on every well-built file.
                faux_b, faux_i = _detect_faux(ps_name, weight >= 600, italic) if ps_name else (False, False)
                family = normalize_font_name(str(style.get("fontFamily") or raw_font))
                if family:
                    doc.fonts.setdefault(family, True)
                page.text.append(
                    TextElement(
                        text=str(node.get("characters", "")),
                        bbox=box,
                        font_family=family,
                        # Figma sizes are px on a 1x frame; 1px == 0.75pt at 96dpi.
                        font_size_pt=float(style.get("fontSize", 16.0)) * 0.75,
                        font_weight=weight,
                        is_bold=weight >= 600,
                        is_italic=italic,
                        is_faux_bold=faux_b,
                        is_faux_italic=faux_i,
                        color_hex=fill_hex,
                        page=page.index,
                        ordinal=counter[0],
                        line_height_pt=float(style["lineHeightPx"]) * 0.75 if style.get("lineHeightPx") else None,
                    )
                )
                counter[0] += 1
            elif box and ntype in ("RECTANGLE", "ELLIPSE", "VECTOR", "STAR", "LINE", "POLYGON", "BOOLEAN_OPERATION"):
                page.shapes.append(
                    ShapeElement(
                        bbox=box,
                        fill_hex=fill_hex,
                        kind=str(ntype),
                        page=page.index,
                        area_norm=(box[2] - box[0]) * (box[3] - box[1]),
                    )
                )
            elif box and (ntype == "IMAGE" or any(f.get("type") == "IMAGE" for f in fills if isinstance(f, dict))):
                page.images.append(ImageElement(bbox=box, page=page.index, name=str(node.get("name", ""))))

            for child in node.get("children", []) or []:
                if isinstance(child, dict):
                    _walk(child)

        _walk(frame)
        doc.pages.append(page)

    if not any(p.text or p.shapes or p.images for p in doc.pages):
        doc.kind = "none"
        doc.warnings.append("Figma frames contained no measurable nodes")
    return doc


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------
def parse_structured_source(
    structured_source: dict[str, Any] | None,
    asset_kind: str,
    raw_bytes: bytes | None = None,
) -> StructuredDocument:
    """Build the element tree from whatever the control plane gave us.

    `structuredSource` on the request wins, because the control plane may have
    parsed a Figma file we cannot reach. Otherwise we parse the bytes ourselves.
    """
    if structured_source:
        kind = str(structured_source.get("kind") or structured_source.get("type") or "").lower()
        if kind in ("figma", "figma-file", "figma_node") or "document" in structured_source:
            figma = parse_figma(structured_source)
            if figma.available:
                return figma
        if kind == "html" and isinstance(structured_source.get("html"), str):
            html = parse_html(structured_source["html"])
            if html.available:
                return html
        if kind in ("pdf", "pptx") and isinstance(structured_source.get("pages"), list):
            adopted = _adopt_prepared(structured_source, kind)  # type: ignore[arg-type]
            if adopted.available:
                return adopted

    if raw_bytes:
        try:
            if asset_kind == "pdf" or raw_bytes[:5] == b"%PDF-":
                return parse_pdf(raw_bytes)
            if asset_kind == "pptx" or raw_bytes[:2] == b"PK":
                return parse_pptx(raw_bytes)
            if asset_kind == "html":
                return parse_html(raw_bytes.decode("utf-8", errors="replace"))
        except Exception as exc:  # noqa: BLE001
            log.warning("structured_parse_failed", kind=asset_kind, error=str(exc))
            return StructuredDocument(kind="none", warnings=[f"structured parse failed: {exc}"])

    return StructuredDocument(kind="none")


def _adopt_prepared(payload: dict[str, Any], kind: SourceKind) -> StructuredDocument:
    """Adopt an element tree the control plane already normalized."""
    doc = StructuredDocument(kind=kind)
    for i, raw_page in enumerate(payload.get("pages", [])):
        if not isinstance(raw_page, dict):
            continue
        page = StructuredPage(
            index=int(raw_page.get("index", i)),
            width_pt=float(raw_page.get("widthPt", raw_page.get("width", 612.0))),
            height_pt=float(raw_page.get("heightPt", raw_page.get("height", 792.0))),
        )
        for j, raw in enumerate(raw_page.get("text", []) or []):
            if not isinstance(raw, dict):
                continue
            bbox = raw.get("bbox") or [0, 0, 0, 0]
            raw_font = str(raw.get("fontFamily", raw.get("font", "")))
            declared_bold = bool(raw.get("bold", raw.get("isBold", False)))
            declared_italic = bool(raw.get("italic", raw.get("isItalic", False)))
            faux_b, faux_i = _detect_faux(raw_font, declared_bold, declared_italic)
            page.text.append(
                TextElement(
                    text=str(raw.get("text", "")),
                    bbox=tuple(float(v) for v in bbox[:4]),  # type: ignore[arg-type]
                    font_family=normalize_font_name(raw_font),
                    font_size_pt=float(raw.get("fontSizePt", raw.get("size", 0.0)) or 0.0),
                    font_weight=int(raw.get("fontWeight", _weight_from_name(raw_font, declared_bold))),
                    is_bold=declared_bold,
                    is_italic=declared_italic,
                    is_faux_bold=faux_b,
                    is_faux_italic=faux_i,
                    color_hex=raw.get("color") or raw.get("colorHex"),
                    page=page.index,
                    ordinal=j,
                )
            )
        for raw in raw_page.get("shapes", []) or []:
            if not isinstance(raw, dict):
                continue
            bbox = raw.get("bbox") or [0, 0, 0, 0]
            b = tuple(float(v) for v in bbox[:4])
            page.shapes.append(
                ShapeElement(
                    bbox=b,  # type: ignore[arg-type]
                    fill_hex=raw.get("fill") or raw.get("fillHex"),
                    stroke_hex=raw.get("stroke"),
                    kind=str(raw.get("kind", "rect")),
                    page=page.index,
                    area_norm=(b[2] - b[0]) * (b[3] - b[1]),
                )
            )
        doc.pages.append(page)
    if not doc.pages:
        doc.kind = "none"
    return doc


__all__ = [
    "ImageElement",
    "ShapeElement",
    "SourceKind",
    "StructuredDocument",
    "StructuredPage",
    "TextElement",
    "normalize_font_name",
    "parse_figma",
    "parse_html",
    "parse_pdf",
    "parse_pptx",
    "parse_structured_source",
]

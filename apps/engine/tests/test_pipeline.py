"""End-to-end: tier ordering, budget degradation, graceful failure, routes."""

from __future__ import annotations

import json

import numpy as np
import pytest

from brandlens_engine import ENGINE_VERSION
from brandlens_engine.assemble import assemble, collect_constraints, crop_loss
from brandlens_engine.channel_spec import check_conformance, resolve_spec
from brandlens_engine.induce import induce_rules
from brandlens_engine.models import (
    AssembleBrief,
    AssembleCandidate,
    AssembleRequest,
    InduceRulesRequest,
    RuleDefinition,
)
from brandlens_engine.pipeline import rule_in_scope, run_analysis
from brandlens_engine.registry import ANALYZERS, TIER_ORDER, effective_tier, registered_names, requires_llm

from .conftest import make_rule

# ---------------------------------------------------------------------------
# registry
# ---------------------------------------------------------------------------
REQUIRED_ANALYZERS = [
    "logo.presence", "logo.clearspace", "logo.min_size", "logo.distortion",
    "logo.recolor", "logo.placement", "logo.occlusion",
    "color.palette_conformance", "color.forbidden", "color.dominance_ratio",
    "typography.approved_family", "typography.min_size", "typography.hierarchy",
    "typography.fallback_font", "typography.casing",
    "layout.safe_zone", "layout.margins", "layout.grid_alignment",
    "layout.element_overlap", "layout.text_density",
    "imagery.style_conformance", "imagery.medium", "imagery.prohibited_subject", "imagery.reuse",
    "copy.banned_terms", "copy.required_terms", "copy.readability",
    "copy.claim_substantiation", "copy.disclaimer_present", "copy.locale_spelling", "copy.cta_allowlist",
    "accessibility.contrast", "accessibility.font_size_floor", "accessibility.alt_text",
    "channel_spec.conformance",
    "vlm.voice_tone", "vlm.mood", "vlm.subject_appropriateness",
    "vlm.overall_judgment", "vlm.rule_adjudication",
]


def test_every_contracted_analyzer_is_registered():
    assert set(REQUIRED_ANALYZERS) <= set(registered_names())


def test_registry_is_callable_and_tier_mapped():
    for name in REQUIRED_ANALYZERS:
        assert callable(ANALYZERS[name])
        assert effective_tier(name, "deterministic") in TIER_ORDER


def test_effective_tier_takes_the_stricter_of_declared_and_actual():
    """A rule that mislabels a VLM check as `cv` must not slip past the guard."""
    assert effective_tier("vlm.mood", "cv") == "vlm"
    assert effective_tier("copy.banned_terms", "vlm") == "vlm"
    assert requires_llm("vlm.overall_judgment") is True
    assert requires_llm("copy.banned_terms") is False


# ---------------------------------------------------------------------------
# scoping
# ---------------------------------------------------------------------------
def test_rule_scope_filtering(context_for, poster_path):
    ctx = context_for(poster_path, market="US", channel="instagram")
    in_scope = make_rule("r", "copy.banned_terms", "copy")
    in_scope.scope.markets = ["US", "CA"]
    assert rule_in_scope(in_scope, ctx)

    out = make_rule("r", "copy.banned_terms", "copy")
    out.scope.markets = ["DE"]
    assert not rule_in_scope(out, ctx)


def test_missing_asset_metadata_keeps_a_rule_in_scope(context_for, poster_path):
    """Skipping a rule because metadata is absent is how compliance gaps happen."""
    ctx = context_for(poster_path)
    rule = make_rule("r", "copy.banned_terms", "copy")
    rule.scope.markets = ["DE"]
    assert rule_in_scope(rule, ctx)


# ---------------------------------------------------------------------------
# full run
# ---------------------------------------------------------------------------
def _full_ruleset() -> list[RuleDefinition]:
    return [
        make_rule("channel", "channel_spec.conformance", "channel_spec", params={"spec": {"minWidth": 600}}),
        make_rule("banned", "copy.banned_terms", "copy", severity="blocker"),
        make_rule("required", "copy.required_terms", "copy", severity="blocker"),
        make_rule("a11y-alt", "accessibility.alt_text", "accessibility"),
        make_rule("a11y-contrast", "accessibility.contrast", "accessibility"),
        make_rule("type-family", "typography.approved_family", "typography"),
        make_rule("logo-presence", "logo.presence", "logo", tier="cv", severity="blocker"),
        make_rule("logo-clearspace", "logo.clearspace", "logo", tier="cv", params={"clearSpaceMultiple": 0.25}),
        make_rule("logo-minsize", "logo.min_size", "logo", tier="cv"),
        make_rule("palette", "color.palette_conformance", "color", tier="cv"),
        make_rule("forbidden", "color.forbidden", "color", tier="cv"),
        make_rule("margins", "layout.margins", "layout", tier="cv", params={"minMarginPct": 1.0}),
        make_rule("density", "layout.text_density", "layout", tier="cv"),
        make_rule("style", "imagery.style_conformance", "imagery", tier="cv"),
        make_rule("mood", "vlm.mood", "imagery", tier="vlm"),
        make_rule("overall", "vlm.overall_judgment", "legal", tier="vlm"),
    ]


def test_full_run_produces_one_result_per_rule(make_request, poster_path):
    rules = _full_ruleset()
    request = make_request(poster_path, rules, copy_fields={"body": "Northgate makes money simple."})
    response = run_analysis(request)

    assert len(response.results) == len(rules)
    assert {r.rule_key for r in response.results} == {r.key for r in rules}
    assert response.engine_version == ENGINE_VERSION
    assert response.duration_ms > 0


def test_every_result_carries_a_real_verdict_with_evidence(make_request, poster_path):
    request = make_request(poster_path, _full_ruleset(), copy_fields={"body": "Northgate makes money simple."})
    response = run_analysis(request)

    valid = {"pass", "fail", "not_applicable", "insufficient_evidence", "abstained"}
    for result in response.results:
        assert result.verdict in valid
        assert result.evidence.observation, f"{result.rule_key} produced no observation"
        assert result.latency_ms is not None
        if result.verdict in ("pass", "fail"):
            assert result.evidence.measured, f"{result.rule_key} passed/failed with no measurement"
            assert result.evidence.threshold, f"{result.rule_key} passed/failed with no threshold"


def test_no_analyzer_returns_a_hardcoded_pass(make_request, poster_path):
    """A pass must always be backed by numbers that were actually computed."""
    request = make_request(poster_path, _full_ruleset(), copy_fields={"body": "Northgate makes money simple."})
    response = run_analysis(request)
    passes = [r for r in response.results if r.verdict == "pass"]
    assert passes, "the clean poster should pass at least some criteria"
    for result in passes:
        assert result.evidence.measured
        assert any(
            isinstance(v, (int, float, list, dict)) for v in (result.evidence.measured or {}).values()
        ), f"{result.rule_key} passed with no numeric evidence"


def test_tiers_run_in_order_and_overall_judgment_runs_last(make_request, poster_path):
    request = make_request(poster_path, _full_ruleset(), copy_fields={"body": "Northgate is simple."})
    response = run_analysis(request)
    keys = [r.rule_key for r in response.results]
    # T0 results are banked before T1, and the holistic judgment is dead last.
    assert keys.index("banned") < keys.index("logo-presence")
    assert keys.index("logo-presence") < keys.index("mood")
    assert keys[-1] == "overall"


def test_t2_criteria_abstain_when_no_provider_is_configured(make_request, poster_path):
    request = make_request(poster_path, _full_ruleset(), copy_fields={"body": "Northgate is simple."})
    response = run_analysis(request)
    vlm_results = [r for r in response.results if r.rule_key in ("mood", "overall")]
    assert all(r.verdict == "insufficient_evidence" for r in vlm_results)
    assert response.cost_usd == 0.0
    assert any("judge unavailable" in w for w in response.warnings)


def test_deterministic_only_skips_t2_and_flags_degraded(make_request, poster_path):
    request = make_request(poster_path, _full_ruleset(), copy_fields={"body": "Northgate is simple."})
    request.deterministic_only = True
    response = run_analysis(request)

    assert response.degraded is True
    assert "deterministicOnly" in (response.degraded_reason or "")
    skipped = [r for r in response.results if r.rule_key in ("mood", "overall")]
    assert all(r.verdict == "insufficient_evidence" for r in skipped)
    assert all(r.evidence.measured.get("reason") == "deterministicOnly" for r in skipped)
    # T0/T1 still produce real answers — a partial answer beats an error.
    assert any(r.verdict in ("pass", "fail") for r in response.results)


def test_budget_ceiling_degrades_instead_of_failing(make_request, poster_path, monkeypatch):
    from brandlens_engine.judge import Judge
    from brandlens_engine.pipeline import AnalysisContext

    from .test_judge import FakeProvider
    from .test_judge import response as judge_response

    rules = _full_ruleset()
    request = make_request(poster_path, rules, copy_fields={"body": "Northgate is simple."})
    request.judge.cost_ceiling_usd = 0.005  # enough for exactly one judge call

    def _judge(self: AnalysisContext) -> Judge:
        if self._judge is None:
            self._judge = Judge(FakeProvider(judge_response("pass", 0.9)), self.judge_config, self.brand)
        return self._judge

    monkeypatch.setattr(AnalysisContext, "judge", _judge)
    result = run_analysis(request)

    assert result.degraded is True
    assert "cost ceiling" in (result.degraded_reason or "")
    assert result.cost_usd > 0
    assert len(result.results) == len(rules)
    skipped = [r for r in result.results if r.evidence.measured and r.evidence.measured.get("reason") == "costCeiling"]
    assert skipped, "the second T2 criterion should have been skipped, not failed"


def test_unknown_analyzer_becomes_insufficient_evidence(make_request, poster_path):
    rules = [make_rule("mystery", "does.not.exist", "copy")]
    response = run_analysis(make_request(poster_path, rules))
    assert response.results[0].verdict == "insufficient_evidence"
    assert "No analyzer is registered" in (response.results[0].evidence.observation or "")
    assert response.results[0].error


def test_a_raising_analyzer_does_not_fail_the_run(make_request, poster_path, monkeypatch):
    import brandlens_engine.registry as registry

    def _boom(ctx, rule):
        raise RuntimeError("simulated analyzer explosion")

    monkeypatch.setitem(registry.ANALYZERS, "copy.banned_terms", _boom)
    rules = [make_rule("boom", "copy.banned_terms", "copy"), make_rule("ok", "channel_spec.conformance", "channel_spec")]
    response = run_analysis(make_request(poster_path, rules, copy_fields={"body": "hi"}))

    boom = next(r for r in response.results if r.rule_key == "boom")
    assert boom.verdict == "insufficient_evidence"
    assert "RuntimeError" in (boom.error or "")
    assert len(response.results) == 2, "the sibling criterion must still be evaluated"


def test_missing_asset_file_degrades_every_pixel_check(make_request):
    rules = _full_ruleset()
    response = run_analysis(make_request("/no/such/asset.png", rules, copy_fields={"body": "Northgate"}))
    assert len(response.results) == len(rules)
    assert not any(r.verdict == "fail" and "hallucinat" in (r.evidence.observation or "") for r in response.results)
    assert any("unavailable" in w or "rasterised" in w for w in response.warnings)
    logo = next(r for r in response.results if r.rule_key == "logo-presence")
    assert logo.verdict == "insufficient_evidence"


def test_out_of_scope_rules_are_reported_not_dropped(make_request, poster_path):
    rule = make_rule("de-only", "copy.banned_terms", "copy")
    rule.scope.markets = ["DE"]
    response = run_analysis(make_request(poster_path, [rule], market="US"))
    assert len(response.results) == 1
    assert response.results[0].verdict == "not_applicable"
    assert "out of scope" in (response.results[0].evidence.observation or "")


def test_measurements_are_exported_and_replayable(make_request, poster_path):
    rules = [make_rule("palette", "color.palette_conformance", "color", tier="cv")]
    first = run_analysis(make_request(poster_path, rules))
    assert first.measurements, "reusable measurements must be returned for the control plane to persist"

    replay = make_request(poster_path, rules)
    replay.cached_measurements = first.measurements
    second = run_analysis(replay)
    assert second.results[0].verdict == first.results[0].verdict


def test_structured_pdf_run_prefers_structure_over_pixels(make_request, brand_pdf_path):
    rules = [
        make_rule("type-family", "typography.approved_family", "typography"),
        make_rule("type-size", "typography.min_size", "typography"),
        make_rule("fallback", "typography.fallback_font", "typography"),
        make_rule("a11y-size", "accessibility.font_size_floor", "accessibility", params={"minSizePt": 8}),
        make_rule("disclaimer", "copy.disclaimer_present", "copy", severity="blocker"),
        make_rule("banned", "copy.banned_terms", "copy", severity="blocker"),
    ]
    response = run_analysis(make_request(brand_pdf_path, rules, kind="pdf"))
    by_key = {r.rule_key: r for r in response.results}

    assert by_key["type-family"].evidence.measured["source"] == "pdf"
    assert by_key["a11y-size"].evidence.measured["exactSizes"] is True
    assert by_key["a11y-size"].verdict == "fail"  # 6.5pt disclaimer
    assert by_key["disclaimer"].verdict == "fail"  # present but below its 8pt floor
    assert by_key["banned"].verdict == "fail"  # "guaranteed" is in the PDF copy
    assert by_key["banned"].evidence.measured["copySource"] == "pdf"


def test_response_serialises_to_the_camelcase_contract(make_request, poster_path):
    response = run_analysis(make_request(poster_path, _full_ruleset(), copy_fields={"body": "Northgate"}))
    payload = json.loads(json.dumps(response.model_dump(by_alias=True, mode="json")))

    assert set(payload) >= {
        "requestId", "results", "measurements", "artifacts", "costUsd",
        "durationMs", "degraded", "degradedReason", "engineVersion", "warnings",
    }
    first = payload["results"][0]
    assert set(first) >= {
        "ruleKey", "ruleVersion", "dimension", "tier", "verdict", "severity",
        "confidence", "evidence", "costUsd", "cached",
    }
    assert "measured" in first["evidence"]


# ---------------------------------------------------------------------------
# channel spec
# ---------------------------------------------------------------------------
def test_channel_spec_resolution_prefers_the_most_specific_key():
    spec = {"instagram:story": {"width": 1080}, "instagram": {"width": 1200}, "_default": {"width": 800}}
    assert resolve_spec(spec, "instagram", "story")[0]["width"] == 1080
    assert resolve_spec(spec, "instagram", "feed")[0]["width"] == 1200
    assert resolve_spec(spec, "tiktok", None)[0]["width"] == 800
    assert resolve_spec(None, "instagram", "story")[0] is None


def test_channel_spec_flags_dimension_and_format_violations(context_for, poster_path, brand):
    brand.channel_spec = {
        "instagram": {"width": 1080, "height": 1920, "allowedFormats": ["jpg"], "maxFileSizeKb": 1}
    }
    ctx = context_for(poster_path, channel="instagram")
    ctx.request.brand = brand
    result = check_conformance(ctx, make_rule("spec", "channel_spec.conformance", "channel_spec"))

    assert result.verdict == "fail"
    constraints = {v["constraint"] for v in result.evidence.measured["violations"]}
    assert {"width", "height", "allowedFormats", "maxFileSizeKb"} <= constraints
    assert result.confidence == 1.0


def test_channel_spec_passes_a_conforming_asset(context_for, poster_path, brand):
    brand.channel_spec = {"_default": {"width": 1200, "height": 628, "allowedFormats": ["png"]}}
    ctx = context_for(poster_path)
    ctx.request.brand = brand
    result = check_conformance(ctx, make_rule("spec", "channel_spec.conformance", "channel_spec"))
    assert result.verdict == "pass"
    assert result.evidence.measured["aspectRatio"] == pytest.approx(1200 / 628, rel=1e-4)


# ---------------------------------------------------------------------------
# induce / assemble
# ---------------------------------------------------------------------------
def test_induction_reports_thresholds_but_proposes_nothing_under_min_support(brand, poster_path):
    from brandlens_engine.models import EngineAssetRef

    assets = [
        EngineAssetRef(id=f"a{i}", kind="image", uri=poster_path, content_hash=f"h{i}") for i in range(4)
    ]
    response = induce_rules(
        InduceRulesRequest(request_id="r", org_id="o", brand_id="b", assets=assets, brand=brand, min_support=20)
    )
    assert response.measured_count == 4
    assert response.rules == []
    assert any("minSupport" in w for w in response.warnings)
    assert response.style_profile is not None
    assert len(response.style_profile["centroid"]) > 0


def test_induced_rules_carry_support_and_stay_proposed(brand, poster_path):
    from brandlens_engine.models import EngineAssetRef

    assets = [
        EngineAssetRef(id=f"a{i}", kind="image", uri=poster_path, content_hash=f"hh{i}") for i in range(8)
    ]
    response = induce_rules(
        InduceRulesRequest(
            request_id="r", org_id="o", brand_id="b", assets=assets, brand=brand, min_support=5, percentile=5
        )
    )
    assert response.rules, "8 assets above a minSupport of 5 should yield proposals"
    for rule in response.rules:
        assert rule.status == "proposed", "induction must never activate a rule"
        assert rule.provenance == "inductive"
        assert rule.support is not None
        assert rule.support.sample_size and rule.support.sample_size >= 5
        assert rule.support.percentile is not None
        assert rule.support.observed_value is not None
        assert rule.check.fn in registered_names()


def test_assemble_enforces_hard_constraints_without_a_model(brand, poster_path):
    brand.channel_spec = {"instagram": {"width": 1080, "height": 1920}}
    rules = [
        make_rule("margins", "layout.margins", "layout", tier="cv", params={"minMarginPct": 8}),
        make_rule("logo", "logo.min_size", "logo", tier="cv", params={"minHeightPct": 9}),
        make_rule("cta", "copy.cta_allowlist", "copy", params={"allowedCtas": ["Open an account"]}),
        make_rule(
            "safe",
            "layout.safe_zone",
            "layout",
            tier="cv",
            params={"zones": [{"name": "caption", "bbox": [0, 0.86, 1, 1]}]},
        ),
    ]
    request = AssembleRequest(
        request_id="r",
        org_id="o",
        brand=brand,
        brief=AssembleBrief(
            title="Spring launch", key_message="Move money the modern way", mandatories=["FDIC insured"]
        ),
        candidate_assets=[
            AssembleCandidate(id="c1", name="Portrait hero", uri=poster_path, width=1080, height=1920, tags=["instagram"]),
            AssembleCandidate(id="c2", name="Wide banner", uri=poster_path, width=1200, height=400, tags=["display"]),
        ],
        rules=rules,
        provider="anthropic",
        model="claude-sonnet-4-5-20250929",
    )
    request.brief.targets = [{"channel": "instagram", "assetType": "story"}]
    response = assemble(request)

    item = response.items[0]
    assert item["widthPx"] == 1080 and item["heightPx"] == 1920
    assert item["backgroundAssetId"] == "c1", "the portrait candidate must beat the wide one"
    assert item["cta"] == "Open an account"
    assert "FDIC insured" in item["mandatoryText"]
    assert item["constraintsEnforced"]["minLogoHeightPct"] == 9
    logo_slot = next(s for s in item["layout"] if s["slot"] == "logo")
    assert logo_slot["bbox"][0] == pytest.approx(0.08), "the 8% margin is honoured by construction"
    legal_slot = next(s for s in item["layout"] if s["slot"] == "legal")
    assert legal_slot["bbox"][3] <= 0.86, "nothing is placed inside the reserved caption zone"
    assert "unavailable" in response.rationale  # no provider configured in tests


def test_collect_constraints_takes_the_strictest_value():
    rules = [
        make_rule("m1", "layout.margins", "layout", params={"minMarginPct": 4}),
        make_rule("m2", "layout.margins", "layout", params={"minMarginPct": 9}),
    ]
    assert collect_constraints(rules, None, {})["minMarginPct"] == 9


def test_crop_loss_math():
    assert crop_loss(1080, 1080, 1080, 1080) == 0.0
    assert crop_loss(1920, 1080, 1080, 1920) == pytest.approx(1 - (1080 / 1920) / (1920 / 1080), abs=1e-3)
    assert crop_loss(0, 0, 10, 10) == 1.0


# ---------------------------------------------------------------------------
# HTTP surface
# ---------------------------------------------------------------------------
def test_health_is_unauthenticated_and_cheap(client):
    r = client.get("/health")
    assert r.status_code == 200
    assert r.json()["status"] == "ok"


def test_version_reports_the_analyzer_count(client):
    body = client.get("/version").json()
    assert body["engineVersion"] == ENGINE_VERSION
    assert body["analyzers"] >= len(REQUIRED_ANALYZERS)
    assert body["python"].startswith("3.")


def test_deep_health_enumerates_analyzers_and_providers(client):
    body = client.get("/health/deep").json()
    assert body["status"] in ("ok", "degraded", "error")
    assert set(body["providers"]) == {"judge", "extract", "text"}
    assert "ocrDriver" in body
    for name in REQUIRED_ANALYZERS:
        assert name in body["analyzers"]


@pytest.mark.parametrize(
    "path",
    ["/v1/analyze", "/v1/extract-rules", "/v1/induce-rules", "/v1/assemble", "/v1/predict", "/v1/embed"],
)
def test_all_non_health_routes_require_the_shared_secret(client, path):
    assert client.post(path, json={}).status_code == 401
    assert client.post(path, json={}, headers={"X-Engine-Secret": "wrong"}).status_code == 401


def test_analyze_route_round_trip(client, auth_headers, make_request, poster_path):
    request = make_request(poster_path, _full_ruleset(), copy_fields={"body": "Northgate is simple."})
    r = client.post(
        "/v1/analyze", json=request.model_dump(by_alias=True, mode="json"), headers=auth_headers
    )
    assert r.status_code == 200
    body = r.json()
    assert body["requestId"] == "req-test"
    assert len(body["results"]) == len(request.rules)
    assert r.headers["x-engine-version"] == ENGINE_VERSION


def test_malformed_body_is_a_422_not_a_500(client, auth_headers):
    r = client.post("/v1/analyze", json={"requestId": "x"}, headers=auth_headers)
    assert r.status_code == 422


def test_embed_route_returns_deterministic_vectors(client, auth_headers):
    payload = {"requestId": "e1", "texts": ["move money the modern way", "something else"], "kind": "text"}
    first = client.post("/v1/embed", json=payload, headers=auth_headers).json()
    second = client.post("/v1/embed", json=payload, headers=auth_headers).json()

    assert first["provider"] == "hash"
    assert first["dim"] == len(first["vectors"][0])
    assert first["vectors"] == second["vectors"], "the hash provider must be reproducible across calls"
    assert first["costUsd"] == 0.0


def test_embed_similarity_is_meaningful_for_near_duplicates(client, auth_headers):
    from brandlens_engine.embeddings import cosine_similarity

    payload = {
        "texts": [
            "move money the modern way",
            "move money the modern way today",
            "quarterly derivatives clearing settlement",
        ],
        "kind": "text",
    }
    vectors = client.post("/v1/embed", json=payload, headers=auth_headers).json()["vectors"]
    near = cosine_similarity(vectors[0], vectors[1])
    far = cosine_similarity(vectors[0], vectors[2])
    assert near > far


def test_predict_without_comparisons_reports_no_percentile(client, auth_headers, poster_path, brand):
    payload = {
        "requestId": "p1",
        "orgId": "o",
        "asset": {"id": "a", "kind": "image", "uri": poster_path, "contentHash": "h"},
        "brand": brand.model_dump(by_alias=True, mode="json"),
        "personas": [{"age": "25-34", "role": "first-time saver"}],
        "comparisonAssets": [],
        "provider": "anthropic",
        "model": "claude-sonnet-4-5-20250929",
    }
    body = client.post("/v1/predict", json=payload, headers=auth_headers).json()
    assert body["percentileVsCorpus"] is None, "an absolute score without references is not publishable"
    assert body["intervalLow"] is None and body["intervalHigh"] is None
    assert any(r["kind"] == "warning" for r in body["recommendations"])


# ---------------------------------------------------------------------------
# structured sources: typography, layout, imagery
# ---------------------------------------------------------------------------
def test_typography_reads_exact_values_from_pptx(context_for, scratch):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    from brandlens_engine.typography import check_approved_family, check_fallback_font, check_min_size

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(6), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Terms apply."
    run.font.size = Pt(7)
    run.font.name = "Comic Sans MS"
    run.font.bold = True
    path = scratch / "deck.pptx"
    prs.save(path)

    ctx = context_for(str(path), kind="pptx")
    family = check_approved_family(ctx, make_rule("t", "typography.approved_family", "typography"))
    assert family.verdict == "fail"
    assert family.evidence.measured["source"] == "pptx"
    assert family.evidence.measured["violations"][0]["font"] == "Comic Sans MS"

    size = check_min_size(ctx, make_rule("t", "typography.min_size", "typography", params={"minSizePt": 9}))
    assert size.verdict == "fail"
    assert size.evidence.measured["smallestPt"] == pytest.approx(7.0, abs=0.01)

    # Bold requested on a family whose name carries no bold cut == faux bold.
    fallback = check_fallback_font(ctx, make_rule("t", "typography.fallback_font", "typography"))
    assert fallback.verdict == "fail"
    assert fallback.evidence.measured["fauxStyling"]


def test_figma_structured_source_is_preferred_over_pixels(context_for, poster_path):
    from brandlens_engine.typography import check_approved_family

    figma = {
        "kind": "figma",
        "document": {
            "type": "FRAME",
            "absoluteBoundingBox": {"x": 0, "y": 0, "width": 1080, "height": 1080},
            "children": [
                {
                    "type": "TEXT",
                    "characters": "Move money the modern way",
                    "absoluteBoundingBox": {"x": 80, "y": 100, "width": 600, "height": 90},
                    "style": {"fontFamily": "Helvetica Neue", "fontSize": 64, "fontWeight": 400},
                    "fills": [{"type": "SOLID", "visible": True, "color": {"r": 0.04, "g": 0.09, "b": 0.2}}],
                }
            ],
        },
    }
    ctx = context_for(poster_path, kind="figma", structured_source=figma)
    result = check_approved_family(ctx, make_rule("t", "typography.approved_family", "typography"))
    assert result.verdict == "pass", "'Helvetica Neue' is a declared alias of the approved family"
    assert result.evidence.measured["source"] == "figma"


def test_typography_casing_and_hierarchy_on_a_pdf(context_for, brand_pdf_path):
    from brandlens_engine.typography import check_casing, check_hierarchy

    ctx = context_for(brand_pdf_path, kind="pdf")
    casing = check_casing(ctx, make_rule("t", "typography.casing", "typography", params={"maxAllCapsRatio": 0.5}))
    assert casing.verdict == "pass"
    assert casing.evidence.measured["allCapsCharRatio"] < 0.5

    hierarchy = check_hierarchy(ctx, make_rule("t", "typography.hierarchy", "typography"))
    assert hierarchy.verdict in ("pass", "fail", "insufficient_evidence")
    assert hierarchy.evidence.observation


def test_layout_analyzers_produce_real_geometry(context_for, poster_path):
    from brandlens_engine.layout import (
        check_element_overlap,
        check_grid_alignment,
        check_safe_zone,
        check_text_density,
    )

    ctx = context_for(poster_path)

    safe = check_safe_zone(
        ctx,
        make_rule("l", "layout.safe_zone", "layout", tier="cv", params={"zones": [{"name": "caption", "bbox": [0, 0.9, 1, 1]}]}),
    )
    assert safe.verdict in ("pass", "fail")
    assert safe.evidence.measured["elementCount"] > 0

    grid = check_grid_alignment(
        ctx, make_rule("l", "layout.grid_alignment", "layout", tier="cv", params={"columns": 12, "marginPct": 4})
    )
    assert grid.evidence.measured["residualPct"]["p50"] >= 0

    overlap = check_element_overlap(ctx, make_rule("l", "layout.element_overlap", "layout", tier="cv"))
    assert overlap.verdict in ("pass", "fail", "not_applicable")

    density = check_text_density(ctx, make_rule("l", "layout.text_density", "layout", tier="cv"))
    assert density.severity == "advisory", "the legacy 20%-text heuristic must never block a release"
    assert density.evidence.threshold["advisoryOnly"] is True


def test_grid_alignment_is_not_applicable_without_a_grid(context_for, poster_path):
    from brandlens_engine.layout import check_grid_alignment

    result = check_grid_alignment(context_for(poster_path), make_rule("l", "layout.grid_alignment", "layout", tier="cv"))
    assert result.verdict == "not_applicable"


def test_imagery_features_and_reuse(context_for, poster_path, scratch):
    from PIL import Image

    from brandlens_engine.imagery import check_medium, check_reuse, extract_style_features, perceptual_hash

    features = extract_style_features(np.asarray(Image.open(poster_path).convert("RGB"), dtype=np.uint8))
    values = features.as_dict()
    assert 0 <= values["lightnessMean"] <= 100
    assert 0 <= values["hueSpread"] <= 1
    assert values["saturationMean"] >= 0

    ctx = context_for(poster_path)
    own_hash = perceptual_hash(np.asarray(Image.open(poster_path).convert("RGB"), dtype=np.uint8))
    reuse = check_reuse(
        ctx, make_rule("i", "imagery.reuse", "imagery", tier="cv", params={"knownHashes": [{"id": "prior", "hash": own_hash}]})
    )
    assert reuse.verdict == "fail", "an identical asset must be caught as reuse"
    assert reuse.evidence.measured["comparisons"][0]["distance"] == 0

    medium = check_medium(ctx, make_rule("i", "imagery.medium", "imagery", tier="cv", params={"allowedMediums": ["photo"]}))
    assert medium.verdict in ("pass", "fail", "insufficient_evidence")
    assert "signals" in medium.evidence.measured


def test_perceptual_hash_survives_rescaling(poster_path):
    from PIL import Image

    from brandlens_engine.imagery import hamming_hex, perceptual_hash

    original = Image.open(poster_path).convert("RGB")
    rescaled = original.resize((original.width // 2, original.height // 2), Image.LANCZOS)
    a = perceptual_hash(np.asarray(original, dtype=np.uint8))
    b = perceptual_hash(np.asarray(rescaled, dtype=np.uint8))
    assert hamming_hex(a, b) <= 8, "the same photo re-exported smaller is the same photo"


def test_style_conformance_needs_an_induced_profile(context_for, poster_path):
    from brandlens_engine.imagery import check_style_conformance

    result = check_style_conformance(context_for(poster_path), make_rule("i", "imagery.style_conformance", "imagery", tier="cv"))
    assert result.verdict == "insufficient_evidence"
    assert "induction" in (result.evidence.observation or "")


# ---------------------------------------------------------------------------
# brand-book extraction
# ---------------------------------------------------------------------------
def test_extraction_chunks_cite_page_and_bbox(scratch):
    import pymupdf

    from brandlens_engine.extract import extract_rules
    from brandlens_engine.models import ExtractRulesRequest

    path = scratch / "brandbook.pdf"
    doc = pymupdf.open()
    page = doc.new_page(width=595, height=842)
    page.insert_text((50, 80), "LOGO CLEAR SPACE", fontname="hebo", fontsize=22)
    page.insert_text((50, 120), "Always maintain clear space of at least 0.5x the logomark height.", fontname="helv", fontsize=11)
    page2 = doc.new_page(width=595, height=842)
    page2.insert_text((50, 80), "COLOUR PALETTE", fontname="hebo", fontsize=22)
    page2.insert_text((50, 120), "Primary blue #0B5FFF. Ink #0A1633. PANTONE 2728 C.", fontname="helv", fontsize=11)
    page2.insert_text((50, 150), "We are direct and plain, we are not jargon-heavy.", fontname="helv", fontsize=11)
    doc.save(path)
    doc.close()

    response = extract_rules(
        ExtractRulesRequest(
            request_id="e1", org_id="o", brand_id="b", document_uri=str(path),
            document_name="Brand Book", mime_type="application/pdf",
            provider="anthropic", model="claude-sonnet-4-5-20250929",
        )
    )

    assert response.page_count == 2
    assert len(response.chunks) >= 2
    # Body copy must survive chunking, not be misread as a section heading.
    assert any("clear space of at least" in c.text for c in response.chunks)
    for chunk in response.chunks:
        assert chunk.page >= 1
        assert chunk.bbox and len(chunk.bbox) == 4
    assert {t.hex for t in response.tokens if t.hex} >= {"#0B5FFF", "#0A1633"}
    assert any(t.type == "spot-color" for t in response.tokens)
    assert response.voice_attributes[0].we_are_not.startswith("jargon")
    # No provider is configured in tests, so proposal is skipped and said so.
    assert response.rules == []
    assert any("ANTHROPIC_API_KEY" in w for w in response.warnings)


def test_extraction_of_an_unreadable_document_degrades(scratch):
    from brandlens_engine.extract import extract_rules
    from brandlens_engine.models import ExtractRulesRequest

    response = extract_rules(
        ExtractRulesRequest(
            request_id="e2", org_id="o", brand_id="b", document_uri=str(scratch / "missing.pdf"),
            document_name="Missing", provider="anthropic", model="m",
        )
    )
    assert response.rules == []
    assert response.page_count == 0
    assert any("could not be read" in w for w in response.warnings)
